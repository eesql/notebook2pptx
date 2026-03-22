import pdfplumber
import re
import io
from pptx import Presentation
from pptx.util import Pt, RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

POINT_TO_EMU = 12700

class UltimatePDF2PPT:
    def __init__(self, pdf_path):
        self.pdf_path = pdf_path
        self.prs = Presentation()
        self.font_map = {'sans': 'Arial', 'inter': 'Segoe UI', 'noto': 'Microsoft YaHei', 'hei': 'Microsoft YaHei'}

    def _convert_coords(self, points):
        return int(points * POINT_TO_EMU)

    def _is_in_table(self, obj, tables):
        """检查一个对象（词或矩形）是否落在已识别的表格区域内"""
        obj_x = (obj['x0'] + obj['x1']) / 2 if 'x0' in obj else obj['x0']
        obj_y = (obj['top'] + obj['bottom']) / 2 if 'top' in obj else obj['top']
        
        for t_bbox in tables:
            x0, top, x1, bottom = t_bbox
            if x0 <= obj_x <= x1 and top <= obj_y <= bottom:
                return True
        return False

    def process(self, output_path):
        with pdfplumber.open(self.pdf_path) as pdf:
            self.prs.slide_width = self._convert_coords(pdf.pages[0].width)
            self.prs.slide_height = self._convert_coords(pdf.pages[0].height)

            for page in pdf.pages:
                slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
                
                # --- 第一步：表格识别 ---
                tables = page.find_tables()
                table_bboxes = [t.bbox for t in tables]
                
                for t in tables:
                    self._add_native_table(slide, t)

                # --- 第二步：矢量背景 (避开表格区域) ---
                for rect in page.rects:
                    if not self._is_in_table(rect, table_bboxes):
                        self._draw_shape(slide, rect)

                # --- 第三步：文本段落 (避开表格区域) ---
                words = page.extract_words(extra_attrs=['size', 'fontname', 'non_stroking_color'])
                # 过滤掉属于表格的文字
                content_words = [w for w in words if not self._is_in_table(w, table_bboxes)]
                paragraphs = self._cluster_paragraphs(content_words)

                for para in paragraphs:
                    self._add_paragraph_to_slide(slide, para)

        self.prs.save(output_path)
        print(f"✨ 10分方案转换完成！")

    def _add_native_table(self, slide, pdf_table):
        """核心：将 PDF 表格转换为 PPT 原生 Table 对象"""
        data = pdf_table.extract()
        rows, cols = len(data), len(data[0])
        
        x0, top, x1, bottom = pdf_table.bbox
        width = x1 - x0
        height = bottom - top

        # 创建 PPT 表格
        shape = slide.shapes.add_table(
            rows, cols, 
            self._convert_coords(x0), self._convert_coords(top), 
            self._convert_coords(width), self._convert_coords(height)
        )
        table = shape.table

        # 填充数据与样式
        for r_idx, row_data in enumerate(data):
            for c_idx, cell_text in enumerate(row_data):
                cell = table.rows[r_idx].cells[c_idx]
                cell.text = str(cell_text) if cell_text else ""
                
                # 默认表格字体设置
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(9) # 默认表格字号
                    paragraph.alignment = PP_ALIGN.CENTER

    def _cluster_paragraphs(self, words, y_tol=5):
        if not words: return []
        words.sort(key=lambda x: (x['top'], x['x0']))
        paras, cur_para = [], [words[0]]
        for i in range(1, len(words)):
            v_gap = words[i]['top'] - words[i-1]['bottom']
            if v_gap < y_tol and words[i]['fontname'] == words[i-1]['fontname']:
                cur_para.append(words[i])
            else:
                paras.append(cur_para); cur_para = [words[i]]
        paras.append(cur_para)
        return paras

    def _add_paragraph_to_slide(self, slide, para_words):
        x0, top = min(w['x0'] for w in para_words), min(w['top'] for w in para_words)
        x1, bottom = max(w['x1'] for w in para_words), max(w['bottom'] for w in para_words)
        
        text = ""
        last_y = para_words[0]['top']
        for w in para_words:
            if w['top'] - last_y > 3: text += "\n"
            text += w['text'] + " "
            last_y = w['top']

        txBox = slide.shapes.add_textbox(self._convert_coords(x0), self._convert_coords(top), 
                                         self._convert_coords(x1-x0), self._convert_coords(bottom-top))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text.strip()
        p.font.size = Pt(para_words[0]['size'])
        # 字体映射逻辑...（省略与之前一致的颜色逻辑）

    def _draw_shape(self, slide, rect):
        try:
            s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._convert_coords(rect['x0']), 
                                      self._convert_coords(rect['top']), self._convert_coords(rect['width']), 
                                      self._convert_coords(rect['height']))
            s.line.fill.background()
            if 'non_stroking_color' in rect:
                c = rect['non_stroking_color']
                if len(c) == 3: s.fill.fore_color.rgb = RGBColor(*(int(x*255) for x in c))
        except: pass

# 运行示例
converter = ProfessionalPDF2PPT("notebook.pdf")
converter.process("output_v7.pptx")